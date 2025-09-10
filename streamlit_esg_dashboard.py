import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
from io import BytesIO
import logging

# -------------- Page Config ----------------
st.set_page_config(
    page_title="ESG Risk Dashboard",
    layout="wide",
    initial_sidebar_state="expanded"
)

# -------------- Logging --------------------
logging.basicConfig(level=logging.INFO)

# -------------- Helpers --------------------
@st.cache_data
def sample_data():
    """Generate sample ESG dataset for demo purposes"""
    companies = ["Acme Corp", "GreenTech Ltd", "BlueEnergy", "LocalFoods"]
    categories = ["Environmental", "Social", "Governance"]
    terms = ["emissions", "fine", "lawsuit", "non-compliance", "pollution", "worker_safety", "data_breach"]
    rows = []
    rng = np.random.default_rng(42)

    weight_map = {"lawsuit": 3, "fine": 2, "non-compliance": 2,
                  "emissions": 1, "pollution": 1, "worker_safety": 2, "data_breach": 3}

    for c in companies:
        for cat in categories:
            for t in rng.choice(terms, size=3, replace=False):
                count = int(rng.integers(0, 10))
                if count == 0:
                    continue
                sentiment = float(np.round(rng.uniform(-1, 0.8), 2))
                doc = f"{rng.choice([2022, 2023, 2024])}_Sustainability"
                rows.append({
                    "company": c,
                    "document": str(doc),
                    "category": cat,
                    "term": t,
                    "count": count,
                    "sentiment": sentiment,
                    "risk_weight": float(weight_map.get(t, 1))
                })
    return pd.DataFrame(rows)


def load_uploaded_files(uploaded_files):
    """Load CSV/XLSX files and concatenate"""
    dfs = []
    for f in uploaded_files:
        try:
            if f.name.lower().endswith(".csv"):
                d = pd.read_csv(f)
            else:
                d = pd.read_excel(f)
            dfs.append(d)
        except Exception as e:
            st.sidebar.error(f"Failed to read {f.name}: {e}")
            logging.error(f"File load error: {e}")
    if dfs:
        return pd.concat(dfs, ignore_index=True)
    return pd.DataFrame()


def coerce_schema(df):
    """Ensure required columns exist, add defaults if missing, remove duplicates"""
    required_cols = ['company', 'document', 'category', 'term', 'count', 'sentiment', 'risk_weight']
    for c in required_cols:
        if c not in df.columns:
            if c == 'count':
                df[c] = 1
            elif c == 'sentiment':
                df[c] = 0.0
            elif c == 'risk_weight':
                df[c] = 1.0
            elif c == 'document':
                df[c] = 'unknown'
            else:
                df[c] = 'Unknown'

    # Force types
    df['company'] = df['company'].astype(str)
    df['document'] = df['document'].astype(str)
    df['category'] = df['category'].astype(str)
    df['term'] = df['term'].astype(str)
    df['count'] = pd.to_numeric(df['count'], errors='coerce').fillna(0).astype(int)
    df['sentiment'] = pd.to_numeric(df['sentiment'], errors='coerce').fillna(0.0)
    df['risk_weight'] = pd.to_numeric(df['risk_weight'], errors='coerce').fillna(1.0)

    # Remove duplicates by aggregating
    df = df.groupby(['company', 'document', 'category', 'term']).agg({
        'count': 'sum',
        'sentiment': 'mean',
        'risk_weight': 'mean'
    }).reset_index()
    return df


@st.cache_data
def compute_aggregates(df):
    """Compute total flags, weighted score, sentiment balance, ESG index"""
    if df.empty:
        return pd.DataFrame(), pd.DataFrame()

    df['weighted'] = df['count'] * df['risk_weight']
    df['sent_weighted_sent'] = df['sentiment'] * df['count']

    totals = df.groupby('company', as_index=False).agg(total_flags=('count', 'sum'))
    weighted = df.groupby('company', as_index=False).agg(weighted_score=('weighted', 'sum'))
    sent = df.groupby('company', as_index=False).agg(sentiment_balance=('sent_weighted_sent', 'sum'))
    counts = df.groupby(['company', 'category'], as_index=False).agg(cat_count=('count', 'sum'))

    agg = totals.merge(weighted, on='company').merge(sent, on='company')
    agg['ESG_Index'] = agg.apply(
        lambda r: (r['weighted_score'] + abs(r['sentiment_balance'])) / r['total_flags']
        if r['total_flags'] > 0 else 0, axis=1
    )
    agg['Percentile'] = agg['ESG_Index'].rank(pct=True)
    return agg, counts


def risk_level(score):
    if score > 1000:
        return "🔴 High"
    elif score > 500:
        return "🟠 Medium"
    else:
        return "🟢 Low"


def download_df_as_csv_bytes(df):
    return df.to_csv(index=False).encode('utf-8')


def extract_year_from_document(s):
    import re
    m = re.search(r'(\b20\d{2}\b)', s)
    return int(m.group(1)) if m else None

# ---------------- Sidebar ------------------
st.sidebar.title("Data & Filters")
uploaded = st.sidebar.file_uploader(
    "Upload CSV/XLSX files (NLP outputs)",
    type=['csv', 'xlsx'],
    accept_multiple_files=True
)

raw_df = load_uploaded_files(uploaded) if uploaded else sample_data()
if raw_df.empty:
    st.sidebar.warning("Using demo data")
    raw_df = sample_data()

df = coerce_schema(raw_df)
df['year'] = df['document'].apply(lambda x: extract_year_from_document(str(x)))
if 'sector' not in df.columns:
    df['sector'] = 'Unknown'

# Filters
all_companies = sorted(df['company'].unique())
selected_companies = st.sidebar.multiselect("Select companies", options=all_companies, default=all_companies)
all_docs = sorted(df['document'].unique())
selected_docs = st.sidebar.multiselect("Select documents", options=all_docs, default=all_docs)
years = sorted([y for y in df['year'].unique() if pd.notna(y)])
selected_years = st.sidebar.multiselect("Select years", options=years, default=years if years else [])
all_categories = sorted(df['category'].unique())
selected_cats = st.sidebar.multiselect("Select categories", options=all_categories, default=all_categories)
min_count = st.sidebar.number_input("Minimum count filter", min_value=0, value=1)

df_filtered = df[
    (df['company'].isin(selected_companies)) &
    (df['document'].isin(selected_docs)) &
    (df['category'].isin(selected_cats)) &
    (df['count'] >= int(min_count)) &
    ((df['year'].isin(selected_years)) if selected_years else True)
].copy()

# ---------------- Main Layout ----------------
st.title("AI-Driven ESG Risk Dashboard — Internship Ready")
st.markdown("Compare ESG flagged risks across companies. Adjust filters in the sidebar.")

tab1, tab2, tab3 = st.tabs(["📊 Comparison", "🔎 Drill-Down", "📥 Downloads"])

# -------- Tab 1: Comparison --------
with tab1:
    st.header("Company Comparison & Benchmarks")
    agg_df, cat_counts = compute_aggregates(df_filtered)
    if agg_df.empty:
        st.info("No data to display. Adjust filters.")
    else:
        agg_df['Risk Level'] = agg_df['weighted_score'].apply(risk_level)

        c1, c2 = st.columns([2, 1])
        with c1:
            st.subheader("Total Flags vs Weighted Score")
            fig = px.bar(
                agg_df.sort_values('weighted_score', ascending=False),
                x='company',
                y=['total_flags', 'weighted_score'],
                barmode='group',
                color_discrete_sequence=['#1f77b4', '#ff7f0e']
            )
            st.plotly_chart(fig, use_container_width=True)

            if df_filtered['year'].notna().any():
                st.subheader("Trend: Total Flags over Years")
                yearly = df_filtered.groupby(['company', 'year'], as_index=False).agg(total_flags=('count', 'sum'))
                fig_trend = px.line(yearly, x='year', y='total_flags', color='company', markers=True)
                st.plotly_chart(fig_trend, use_container_width=True)

        with c2:
            st.subheader("Company Scorecards")
            display_df = agg_df[['company', 'total_flags', 'weighted_score', 'ESG_Index', 'Risk Level', 'sentiment_balance']].copy()
            display_df = display_df.rename(columns={
                'total_flags': 'Total Flags',
                'weighted_score': 'Weighted Score',
                'ESG_Index': 'ESG Index',
                'sentiment_balance': 'Sentiment Balance'
            })
            st.dataframe(
                display_df.style.format({'Weighted Score': '{:.1f}', 'ESG Index': '{:.2f}', 'Sentiment Balance': '{:.2f}'}),
                height=420
            )

        # Heatmap
        if not cat_counts.empty:
            st.subheader("Category Distribution Heatmap")
            pivot = cat_counts.pivot(index='company', columns='category', values='cat_count').fillna(0)
            fig2 = px.imshow(
                pivot,
                labels=dict(x="Category", y="Company", color="Count"),
                aspect='auto',
                color_continuous_scale='Viridis'
            )
            st.plotly_chart(fig2, use_container_width=True)

        # Radar chart
        radar_df = cat_counts.pivot(index='company', columns='category', values='cat_count').fillna(0).reset_index()
        if radar_df.shape[0] >= 1:
            st.subheader("ESG Balance Radar Chart")
            melt = radar_df.melt(id_vars='company', var_name='category', value_name='count')
            fig_radar = px.line_polar(
                melt, r='count', theta='category', color='company',
                line_close=True, title='ESG Category Balance (Radar)'
            )
            st.plotly_chart(fig_radar, use_container_width=True)

# -------- Tab 2: Drill-Down --------
with tab2:
    st.header("Drill-Down Insights")
    companies_list = sorted(df_filtered['company'].unique())
    selected = st.selectbox("Choose company to inspect", options=companies_list if companies_list else ["None"])
    if selected != "None":
        comp_df = df_filtered[df_filtered['company'] == selected].copy()
        st.subheader(f"Top flagged terms for {selected}")
        top_n = st.slider("Number of top terms", min_value=3, max_value=20, value=10)
        top_terms = comp_df.groupby('term', as_index=False).agg(total=('count', 'sum')).sort_values('total', ascending=False)
        st.table(top_terms.head(top_n).reset_index(drop=True))

        st.subheader("Sentiment Distribution")
        comp_df['sent_bin'] = pd.cut(comp_df['sentiment'], bins=[-1, -0.25, 0, 0.25, 1],
                                     labels=['Very Negative', 'Negative', 'Positive', 'Very Positive'])
        sent_df = comp_df.groupby('sent_bin').agg(weighted_count=('count', 'sum')).reset_index()
        fig_sent = px.bar(sent_df, x='sent_bin', y='weighted_count', title=f'Sentiment Buckets for {selected}')
        st.plotly_chart(fig_sent, use_container_width=True)

        st.subheader("Detailed Flagged Rows")
        st.dataframe(
            comp_df[['document', 'category', 'term', 'count', 'sentiment', 'risk_weight']]
            .sort_values('count', ascending=False).reset_index(drop=True)
        )

# -------- Tab 3: Downloads & Export --------
with tab3:
    st.header("Downloads & Export")
    if df_filtered.empty:
        st.info("No data to download. Adjust filters.")
    else:
        # CSV downloads
        st.download_button(
            "⬇ Download Filtered CSV",
            data=download_df_as_csv_bytes(df_filtered),
            file_name="esg_filtered.csv",
            mime="text/csv"
        )
        agg_df, counts_df = compute_aggregates(df_filtered)
        if not agg_df.empty:
            st.download_button(
                "⬇ Download Company Scorecards CSV",
                data=download_df_as_csv_bytes(agg_df),
                file_name="esg_scorecards.csv",
                mime="text/csv"
            )

        # PPT Export
        st.markdown("### 📑 PowerPoint Export")
        st.caption("Creates a PPTX with summary tables and charts for your presentation.")
        if st.button("Generate PPTX"):
            try:
                import plotly.io as pio
                from pptx import Presentation
                from pptx.util import Inches

                prs = Presentation()
                slide_layout = prs.slide_layouts[5]  # blank layout

                # Slide 1: Summary Table
                slide1 = prs.slides.add_slide(slide_layout)
                slide1.shapes.title.text = "ESG Risk Dashboard - Summary"
                rows, cols = agg_df.shape
                table = slide1.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(2)).table
                for j, col in enumerate(agg_df.columns):
                    table.cell(0, j).text = str(col)
                for i, row in agg_df.iterrows():
                    for j, col in enumerate(agg_df.columns):
                        table.cell(i+1, j).text = str(row[col])

                # Slide 2: Top flagged terms
                slide2 = prs.slides.add_slide(slide_layout)
                slide2.shapes.title.text = "Top Flagged Terms"
                top_offset = 1.5
                for idx, comp in enumerate(agg_df['company']):
                    comp_df = df_filtered[df_filtered['company'] == comp]
                    top_terms = comp_df.groupby('term', as_index=False).agg(total=('count', 'sum')).sort_values('total', ascending=False).head(5)
                    top_str = "\n".join([f"{r.term}: {r.total}" for r in top_terms.itertuples()])
                    textbox = slide2.shapes.add_textbox(Inches(0.5), Inches(top_offset + idx*1.2), Inches(4), Inches(1))
                    textbox.text = f"{comp}:\n{top_str}"

                # Slide 3: Category heatmap
                if not counts_df.empty:
                    slide3 = prs.slides.add_slide(slide_layout)
                    slide3.shapes.title.text = "Category Heatmap"
                    pivot = counts_df.pivot(index='company', columns='category', values='cat_count').fillna(0)
                    fig_heat = px.imshow(pivot, labels=dict(x="Category", y="Company", color="Count"),
                                         aspect='auto', color_continuous_scale='Viridis')
                    img_bytes = pio.to_image(fig_heat, format='png', width=1200, height=600, scale=2)
                    slide3.shapes.add_picture(BytesIO(img_bytes), Inches(0.5), Inches(1.5), width=Inches(9))

                # Slide 4: Trend chart
                if 'year' in df_filtered.columns and df_filtered['year'].notna().any():
                    slide4 = prs.slides.add_slide(slide_layout)
                    slide4.shapes.title.text = "Yearly Trend - Total Flags"
                    yearly = df_filtered.groupby(['company', 'year'], as_index=False).agg(total_flags=('count', 'sum'))
                    fig_trend = px.line(yearly, x='year', y='total_flags', color='company', markers=True)
                    img_bytes = pio.to_image(fig_trend, format='png', width=1200, height=600, scale=2)
                    slide4.shapes.add_picture(BytesIO(img_bytes), Inches(0.5), Inches(1.5), width=Inches(9))

                bio = BytesIO()
                prs.save(bio)
                bio.seek(0)
                st.download_button(
                    "⬇ Download PPTX",
                    data=bio,
                    file_name="esg_dashboard.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"PPTX generation failed: {e}")
